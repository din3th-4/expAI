"""Continue the EXP-AI presentation with the generated findings and graphs."""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path("presentation/AI_LAB_expai_presentation_ACSAI_2025_2026_titles_fixed.pptx")
OUT = Path("presentation/AI_LAB_expai_presentation_ACSAI_2025_2026_completed.pptx")
GRAPHS = Path("study_outputs/graphs")
DATA = Path("study_outputs")

BLUE = RGBColor(31, 78, 121)
BLACK = RGBColor(35, 35, 35)
GREY = RGBColor(85, 85, 85)
ORANGE = RGBColor(242, 142, 43)


def delete_shape(shape):
    shape.element.getparent().remove(shape.element)


def clear_slide(slide):
    for shape in list(slide.shapes):
        delete_shape(shape)


def add_footer(slide, number, total):
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(7.46), Inches(16), Inches(0.54)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(242, 242, 242)
    bar.line.fill.background()

    num = slide.shapes.add_textbox(Inches(0.12), Inches(7.55), Inches(0.9), Inches(0.25))
    tf = num.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{number}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = GREY

    foot = slide.shapes.add_textbox(Inches(7.0), Inches(7.54), Inches(8.5), Inches(0.3))
    tf = foot.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "AI Lab: Human and AI-Based Emotion Recognition Under Image Degradation"
    p.font.size = Pt(10)
    p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, title, eyebrow=None):
    if eyebrow:
        box = slide.shapes.add_textbox(Inches(2.0), Inches(0.18), Inches(12.2), Inches(0.32))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = eyebrow
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ORANGE

    box = slide.shapes.add_textbox(Inches(2.0), Inches(0.58), Inches(13.2), Inches(0.58))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = BLUE


def add_bullets(slide, bullets, left, top, width, height, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)
    return box


def add_note(slide, text, left, top, width, height, size=13):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.italic = True
    p.font.color.rgb = GREY
    return box


def add_image(slide, filename, left, top, width=None, height=None):
    path = GRAPHS / filename
    if not path.exists():
        raise FileNotFoundError(path)
    kwargs = {"left": Inches(left), "top": Inches(top)}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    return slide.shapes.add_picture(str(path), **kwargs)


def add_logo_placeholder(slide):
    # Keep the top-left visual balance used by the template without depending on
    # the original embedded image object.
    box = slide.shapes.add_textbox(Inches(0.42), Inches(0.28), Inches(1.2), Inches(0.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "EXP-AI"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE


def setup_slide(slide, number, total, title, eyebrow=None):
    clear_slide(slide)
    add_logo_placeholder(slide)
    add_title(slide, title, eyebrow)
    add_footer(slide, number, total)


def fix_existing_bounds(prs):
    """Nudge pre-existing out-of-bounds template objects back onto the canvas."""
    for slide_number in [2, 5]:
        slide = prs.slides[slide_number - 1]
        for shape in slide.shapes:
            if shape.top < 0:
                shape.top = 0

    # Slide 10 has an inherited text box extending beyond the right edge.
    slide = prs.slides[9]
    for shape in slide.shapes:
        if shape.left + shape.width > prs.slide_width:
            shape.width = prs.slide_width - shape.left - Inches(0.25)


def numbers():
    condition = pd.read_csv(DATA / "condition_summary.csv")
    comparison = pd.read_csv(DATA / "human_ai_comparison.csv")
    robustness = pd.read_csv(DATA / "modification_robustness.csv")
    models = pd.read_csv(DATA / "model_emotion_profiles.csv")
    human_version = pd.read_csv(DATA / "human_version_profiles.csv")

    mean_agreement = comparison.groupby("source")[["cosine_similarity"]].mean()
    dominant = comparison.groupby("source")["dominant_agreement"].mean()
    disgust_mean = models[models["detected"].astype(str).str.lower().isin(["true", "1"])]
    disgust_mean = disgust_mean.groupby("source")["disgust"].mean()
    disgust_dom = pd.crosstab(models["source"], models["dominant_emotion"])
    robust = robustness.pivot_table(
        index="version_name", columns="source", values="cosine_to_original"
    )
    human_dom = pd.crosstab(human_version["version_name"], human_version["dominant_emotion"])

    return {
        "fer_cosine": mean_agreement.loc["FER", "cosine_similarity"],
        "df_cosine": mean_agreement.loc["DeepFace", "cosine_similarity"],
        "fer_dom": dominant.loc["FER"],
        "df_dom": dominant.loc["DeepFace"],
        "fer_disgust_mean": disgust_mean.loc["FER"],
        "df_disgust_mean": disgust_mean.loc["DeepFace"],
        "fer_disgust_dom": int(disgust_dom.get("disgust", pd.Series()).get("FER", 0)),
        "df_disgust_dom": int(disgust_dom.get("disgust", pd.Series()).get("DeepFace", 0)),
        "df_blur_robust": robust.loc["Blur", "DeepFace"],
        "df_grey_robust": robust.loc["Greyscale", "DeepFace"],
        "human_grey_sad": int(human_dom.loc["Greyscale"].get("sad", 0)),
        "human_grey_neutral": int(human_dom.loc["Greyscale"].get("neutral", 0)),
        "human_grey_disgust": int(human_dom.loc["Greyscale"].get("disgust", 0)),
    }


def main():
    if not BASE.exists():
        raise FileNotFoundError(BASE)
    prs = Presentation(str(BASE))
    total = len(prs.slides)
    n = numbers()
    fix_existing_bounds(prs)

    # Slide 11: roadmap into results.
    setup_slide(
        prs.slides[10],
        11,
        total,
        "Results roadmap",
        "From data coverage to interpretation",
    )
    add_bullets(
        prs.slides[10],
        [
            "First, check whether the data sources are complete enough to compare.",
            "Then compare human pooled responses with FER and DeepFace outputs.",
            "Next, inspect where dominant emotions collide or disagree.",
            "Finally, test modified images and connect results to MediaPipe facial features.",
        ],
        2.1,
        1.7,
        11.8,
        3.8,
        21,
    )
    add_note(
        prs.slides[10],
        "Important framing: because the images are web-selected, the deck discusses agreement and disagreement, not objective correctness.",
        2.1,
        6.1,
        11.8,
        0.7,
        15,
    )

    # Slide 12: coverage.
    setup_slide(prs.slides[11], 12, total, "The dataset is usable, but not equally complete", "Data coverage")
    add_image(prs.slides[11], "01_data_coverage.png", 0.75, 1.45, width=7.0)
    add_bullets(
        prs.slides[11],
        [
            "MediaPipe and DeepFace produced usable outputs for all 140 image versions.",
            "FER detected 130 of 140 current images, so FER comparisons use fewer rows.",
            "Human data contains 255 current responses after excluding stale rows.",
        ],
        8.2,
        2.0,
        6.5,
        3.2,
        18,
    )
    add_note(prs.slides[11], "Coverage matters because missing detections can affect model comparison.", 8.2, 5.6, 6.3, 0.6)

    # Slide 13: agreement not accuracy.
    setup_slide(prs.slides[12], 13, total, "Agreement replaces accuracy in this study", "Model-human comparison")
    add_bullets(
        prs.slides[12],
        [
            f"FER had higher average similarity with human pooled scores: {n['fer_cosine']:.2f} vs DeepFace {n['df_cosine']:.2f}.",
            f"Dominant-emotion agreement was also higher for FER: {n['fer_dom']:.2f} vs DeepFace {n['df_dom']:.2f}.",
            "This does not mean FER is objectively “right”; it means FER aligned more with this participant group.",
            "The strongest discussion point is where humans and models converge or diverge.",
        ],
        2.0,
        1.55,
        12.0,
        4.8,
        21,
    )

    # Slide 14: average profiles.
    setup_slide(prs.slides[13], 14, total, "Average emotion profiles show model bias patterns", "Emotion distributions")
    add_image(prs.slides[13], "03_average_emotion_profiles.png", 0.75, 1.35, width=7.6)
    add_bullets(
        prs.slides[13],
        [
            "Humans spread scores across multiple emotions more often than the AI models.",
            f"Both AI models gave very low disgust scores on average: FER {n['fer_disgust_mean']:.3f}, DeepFace {n['df_disgust_mean']:.3f}.",
            "This supports the point that the AI models struggled to represent disgust in this dataset.",
        ],
        8.65,
        1.95,
        5.9,
        3.8,
        18,
    )

    # Slide 15: human heatmap.
    setup_slide(prs.slides[14], 15, total, "Human ratings are multi-dimensional, not single labels", "Human pooled responses")
    add_image(prs.slides[14], "04_human_profile_heatmap.png", 0.7, 1.25, width=7.2)
    add_bullets(
        prs.slides[14],
        [
            "Each image receives a seven-emotion profile instead of one forced label.",
            "Ambiguous images show more spread across emotion categories.",
            "This helps avoid treating researcher labels as ground truth.",
        ],
        8.3,
        2.0,
        6.3,
        3.5,
        18,
    )

    # Slide 16: collisions.
    setup_slide(prs.slides[15], 16, total, "Dominant-emotion matrices show where humans and models collide", "Human vs AI collisions")
    add_image(prs.slides[15], "10_dominant_emotion_matrices.png", 0.6, 1.35, width=9.0)
    add_bullets(
        prs.slides[15],
        [
            f"FER never selected disgust as the dominant emotion in the current outputs; DeepFace selected it only {n['df_disgust_dom']} times.",
            "Disagreement is not random: some human categories are repeatedly mapped to different model categories.",
            "These collisions are more defensible to discuss than “wrong” predictions.",
        ],
        10.0,
        1.95,
        5.0,
        4.0,
        16,
    )

    # Slide 17: agreement by category.
    setup_slide(prs.slides[16], 17, total, "Agreement varies by intended emotion category", "Category-level comparison")
    add_image(prs.slides[16], "11_agreement_by_assigned_category.png", 0.85, 1.35, width=8.1)
    add_bullets(
        prs.slides[16],
        [
            "Some categories are easier for humans and AI to align on than others.",
            "The assigned category is only a grouping variable, not a certified label.",
            "Low agreement categories are useful places to discuss ambiguity.",
        ],
        9.3,
        2.0,
        5.5,
        3.2,
        17,
    )

    # Slide 18: modified section divider.
    setup_slide(prs.slides[17], 18, total, "Modified images test stability under visual change", "Modified image section")
    add_bullets(
        prs.slides[17],
        [
            "Original images were compared with blurred, greyscale, and low-resolution versions.",
            "The question is whether the emotion profile stays similar when visual information changes.",
            "This section compares robustness for humans, FER, and DeepFace.",
        ],
        2.0,
        1.9,
        11.8,
        3.6,
        22,
    )

    # Slide 19: robustness.
    setup_slide(prs.slides[18], 19, total, "Greyscale barely moves the AI models, but blur hurts DeepFace", "Modification robustness")
    add_image(prs.slides[18], "08_modification_robustness_full.png", 0.75, 1.35, width=7.7)
    add_bullets(
        prs.slides[18],
        [
            "Greyscale is almost perfectly stable for FER and DeepFace, likely because model preprocessing reduces color dependence.",
            f"DeepFace has lower robustness on blur ({n['df_blur_robust']:.2f}) than on greyscale ({n['df_grey_robust']:.2f}).",
            "Human ratings move more than the AI models across modified versions.",
        ],
        8.75,
        1.9,
        5.8,
        3.8,
        17,
    )

    # Slide 20: variants and greyscale/ambiguity wording.
    setup_slide(prs.slides[19], 20, total, "Greyscale makes human responses more mixed, not simply more accurate", "Variant interpretation")
    add_bullets(
        prs.slides[19],
        [
            "In greyscale images, human dominant responses cluster around sad, neutral, and disgust.",
            f"In the current data, greyscale produced {n['human_grey_sad']} sad, {n['human_grey_neutral']} neutral, and {n['human_grey_disgust']} disgust dominant human profiles.",
            "So the safer claim is that greyscale increases interpretive uncertainty around darker/negative expressions.",
            "This is exactly why the study uses emotion distributions instead of one fixed answer.",
        ],
        2.0,
        1.55,
        12.2,
        5.1,
        21,
    )

    # Slide 21: ambiguity.
    setup_slide(prs.slides[20], 21, total, "Ambiguous images lower agreement and increase uncertainty", "Ambiguity")
    add_image(prs.slides[20], "07_ambiguity_panels.png", 0.6, 1.35, width=8.4)
    add_bullets(
        prs.slides[20],
        [
            "Ambiguous images produce higher human entropy, meaning responses are more spread out.",
            "Dominant-emotion agreement drops on ambiguous images, especially for DeepFace.",
            "This supports the main framing: the dataset is useful for studying disagreement.",
        ],
        9.4,
        2.0,
        5.4,
        3.5,
        17,
    )

    # Slide 22: MediaPipe divider.
    setup_slide(prs.slides[21], 22, total, "MediaPipe connects emotion scores to facial movements", "MediaPipe section")
    add_bullets(
        prs.slides[21],
        [
            "MediaPipe was used to extract facial blendshape features from each image.",
            "These features are not emotion labels by themselves.",
            "Instead, we correlate them with human, FER, and DeepFace emotion-score profiles.",
        ],
        2.0,
        1.9,
        11.8,
        3.6,
        22,
    )

    # Slide 23: MediaPipe human.
    setup_slide(prs.slides[22], 23, total, "Human ratings connect strongly to intuitive facial features", "MediaPipe and humans")
    add_image(prs.slides[22], "13a_mediapipe_top_associations_human_pooled.png", 0.55, 1.28, width=8.3)
    add_bullets(
        prs.slides[22],
        [
            "Eye widening is strongly associated with human surprise scores.",
            "Mouth smile features are associated with human happiness scores.",
            "Negative correlations are also meaningful: for example, squinting can reduce surprise association.",
        ],
        9.25,
        1.85,
        5.6,
        3.8,
        16,
    )

    # Slide 24: MediaPipe FER.
    setup_slide(prs.slides[23], 24, total, "FER associations emphasize classic expression cues", "MediaPipe and FER")
    add_image(prs.slides[23], "13b_mediapipe_top_associations_fer.png", 0.55, 1.28, width=8.3)
    add_bullets(
        prs.slides[23],
        [
            "FER surprise is strongly associated with eye widening and jaw opening.",
            "FER happiness is strongly associated with left and right mouth-smile features.",
            "This makes FER’s outputs relatively interpretable through facial landmarks.",
        ],
        9.25,
        1.85,
        5.6,
        3.8,
        16,
    )

    # Slide 25: MediaPipe DeepFace + conclusion.
    setup_slide(prs.slides[24], 25, total, "DeepFace shows associations, but aligns less with humans overall", "MediaPipe and conclusion")
    add_image(prs.slides[24], "13c_mediapipe_top_associations_deepface.png", 0.55, 1.28, width=7.8)
    add_bullets(
        prs.slides[24],
        [
            "DeepFace still uses visible facial cues, such as jaw opening for surprise.",
            "However, its human similarity is lower than FER in this dataset.",
            "Overall, the study shows where humans and AI converge, diverge, and react differently to image degradation.",
        ],
        8.85,
        1.65,
        6.0,
        4.4,
        16,
    )

    # Slide 26: thank you, add final project framing.
    clear_slide(prs.slides[25])
    box = prs.slides[25].shapes.add_textbox(Inches(1.4), Inches(2.25), Inches(13.2), Inches(1.0))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Thank you"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER
    sub = prs.slides[25].shapes.add_textbox(Inches(2.0), Inches(3.45), Inches(12.0), Inches(0.8))
    tf = sub.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(34)
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER
    add_footer(prs.slides[25], 26, total)

    prs.save(str(OUT))
    print(f"Saved {OUT.resolve()}")


if __name__ == "__main__":
    main()
