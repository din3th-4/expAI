"""Complete the presentation while preserving the original template styling.

Unlike ``complete_presentation.py``, this version keeps the template chrome:
Sapienza logo, footer bar, footer font/position, slide number, and title boxes.
It only replaces body placeholders and updates slide titles/subtitles.
"""

from __future__ import annotations

from pathlib import Path
import copy

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


BASE = Path("presentation/AI_LAB_expai_presentation_ACSAI_2025_2026_titles_fixed.pptx")
OUT = Path("presentation/AI_LAB_expai_presentation_ACSAI_2025_2026_completed_template_style.pptx")
GRAPHS = Path("study_outputs/graphs")
DATA = Path("study_outputs")

BLACK = RGBColor(30, 30, 30)
GREY = RGBColor(80, 80, 80)


def delete_shape(shape):
    shape.element.getparent().remove(shape.element)


def set_text(shape, text):
    # Preserve the existing PowerPoint run formatting; only replace the visible
    # text. Using ``shape.text = ...`` resets the run properties, which is what
    # caused the title/section labels to lose the template style.
    text_nodes = shape.element.xpath(".//a:t")
    if not text_nodes:
        shape.text = text
        return
    text_nodes[0].text = text
    for node in text_nodes[1:]:
        node.text = ""


def set_title(slide, title, section=None):
    # In this template, shape 4 is the large slide title and shape 5 is the
    # smaller section/progress label.
    if len(slide.shapes) > 4:
        set_text(slide.shapes[4], title)
    if section and len(slide.shapes) > 5:
        set_text(slide.shapes[5], section)


def clear_body_keep_template(slide):
    # Keep shapes 0-5: footer bar, slide number, footer text, logo, title, section.
    for shape in list(slide.shapes)[6:]:
        delete_shape(shape)


def clear_all_shapes(slide):
    for shape in list(slide.shapes):
        delete_shape(shape)


def clone_shape_to_slide(shape, slide):
    new_element = copy.deepcopy(shape.element)
    slide.shapes._spTree.insert_element_before(new_element, "p:extLst")


def apply_content_chrome(slide, template_slide, slide_number):
    """Replace a TOC/blank slide's chrome with normal content-slide chrome."""
    clear_all_shapes(slide)
    for shape in list(template_slide.shapes)[:6]:
        clone_shape_to_slide(shape, slide)
    # The cloned slide number may come from the template slide; update it.
    if len(slide.shapes) > 1:
        slide.shapes[1].text = f"{slide_number}/26"


def add_bullets(slide, bullets, left, top, width, height, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.line_spacing = 1.0
        p.space_before = Pt(11.91)
        p.space_after = Pt(9.92)
    return box


def add_small_note(slide, text, left, top, width, height):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = GREY
    return box


def add_image(slide, filename, left, top, width=None, height=None):
    path = GRAPHS / filename
    if not path.exists():
        raise FileNotFoundError(path)
    kwargs = {"left": Inches(left), "top": Inches(top)}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    return slide.shapes.add_picture(str(path), **kwargs)


def set_toc(slide, text):
    # TOC slides usually have just title + one content shape. Keep the template
    # layout and change only the content text.
    for shape in slide.shapes:
        if getattr(shape, "text", "").strip().startswith("Introduction"):
            shape.text = text
            break


def findings_numbers():
    comparison = pd.read_csv(DATA / "human_ai_comparison.csv")
    robustness = pd.read_csv(DATA / "modification_robustness.csv")
    models = pd.read_csv(DATA / "model_emotion_profiles.csv")
    human_version = pd.read_csv(DATA / "human_version_profiles.csv")

    mean_similarity = comparison.groupby("source")["cosine_similarity"].mean()
    dominant = comparison.groupby("source")["dominant_agreement"].mean()
    detected = models[models["detected"].astype(str).str.lower().isin(["true", "1"])]
    disgust_mean = detected.groupby("source")["disgust"].mean()
    disgust_dom = pd.crosstab(models["source"], models["dominant_emotion"])
    robust = robustness.pivot_table(
        index="version_name", columns="source", values="cosine_to_original"
    )
    human_dom = pd.crosstab(human_version["version_name"], human_version["dominant_emotion"])

    return {
        "fer_similarity": mean_similarity.loc["FER"],
        "df_similarity": mean_similarity.loc["DeepFace"],
        "fer_dom": dominant.loc["FER"],
        "df_dom": dominant.loc["DeepFace"],
        "fer_disgust_mean": disgust_mean.loc["FER"],
        "df_disgust_mean": disgust_mean.loc["DeepFace"],
        "fer_disgust_dom": int(disgust_dom.get("disgust", pd.Series()).get("FER", 0)),
        "df_disgust_dom": int(disgust_dom.get("disgust", pd.Series()).get("DeepFace", 0)),
        "df_blur": robust.loc["Blur", "DeepFace"],
        "df_grey": robust.loc["Greyscale", "DeepFace"],
        "human_grey_sad": int(human_dom.loc["Greyscale"].get("sad", 0)),
        "human_grey_neutral": int(human_dom.loc["Greyscale"].get("neutral", 0)),
        "human_grey_disgust": int(human_dom.loc["Greyscale"].get("disgust", 0)),
    }


def fill_two_column_slide(slide, title, section, image_file, bullets):
    clear_body_keep_template(slide)
    set_title(slide, title, section)
    add_image(slide, image_file, 0.55, 1.45, width=7.5)
    add_bullets(slide, bullets, 8.45, 1.8, 6.45, 4.6, size=16)


def fill_text_slide(slide, title, section, bullets, note=None):
    clear_body_keep_template(slide)
    set_title(slide, title, section)
    add_bullets(slide, bullets, 1.9, 1.7, 12.2, 4.6, size=20)
    if note:
        add_small_note(slide, note, 1.9, 6.35, 12.2, 0.55)


def main():
    if not BASE.exists():
        raise FileNotFoundError(BASE)
    prs = Presentation(str(BASE))
    n = findings_numbers()
    chrome_template = prs.slides[11]

    # Keep the user's earlier slides, but fix the one inherited body box that
    # extends past the right edge.
    slide_10 = prs.slides[9]
    for shape in slide_10.shapes:
        if shape.left + shape.width > prs.slide_width:
            shape.width = prs.slide_width - shape.left - Inches(0.25)

    # Keep slides 1-10 as the user's drafted intro/method slides. Convert later
    # TOC slides into normal content slides so footer/logo/title styling stays
    # consistent.
    for slide_number in [11, 14, 18, 22, 24]:
        apply_content_chrome(prs.slides[slide_number - 1], chrome_template, slide_number)

    fill_text_slide(
        prs.slides[10],
        "Results roadmap",
        "Results overview ●",
        [
            "Start with data coverage to check which sources are complete.",
            "Compare humans, FER, and DeepFace using agreement rather than certified accuracy.",
            "Show where emotions collide, then move into image variants and ambiguity.",
            "Finish with MediaPipe associations for humans, FER, and DeepFace.",
        ],
        "This keeps the presentation focused on comparison and interpretation, not ground-truth correctness.",
    )

    fill_text_slide(
        prs.slides[11],
        "The dataset is usable, but not equally complete",
        "Results ●○○",
        [
            "DeepFace and MediaPipe produced usable outputs for all 140 image versions.",
            "FER detected 130 of 140 current images, so some FER comparisons use fewer rows.",
            "Human responses were filtered to 255 current ratings after stale rows were excluded.",
            "This coverage check is important before comparing humans and AI models.",
        ],
        "Main framing: because the dataset is web-selected, we discuss agreement rather than objective accuracy.",
    )
    add_image(prs.slides[11], "01_data_coverage.png", 7.9, 1.55, width=6.15)

    fill_text_slide(
        prs.slides[12],
        "Agreement replaces accuracy in this study",
        "Results ●○○",
        [
            f"FER aligned more closely with human pooled ratings: cosine similarity {n['fer_similarity']:.2f}.",
            f"DeepFace aligned less closely with human pooled ratings: cosine similarity {n['df_similarity']:.2f}.",
            f"Dominant-emotion agreement was also higher for FER ({n['fer_dom']:.2f}) than DeepFace ({n['df_dom']:.2f}).",
            "This does not mean FER is objectively correct; it means FER behaved more similarly to this participant group.",
        ],
    )

    fill_two_column_slide(
        prs.slides[13],
        "Average emotion profiles reveal model bias patterns",
        "Results ●○○",
        "03_average_emotion_profiles.png",
        [
            "Humans spread emotion scores across categories more than the AI models.",
            f"Both AI models gave very low disgust scores on average: FER {n['fer_disgust_mean']:.3f}, DeepFace {n['df_disgust_mean']:.3f}.",
            "This supports the point that the models struggled to represent disgust in this dataset.",
        ],
    )

    fill_two_column_slide(
        prs.slides[14],
        "Human ratings are multi-dimensional, not single labels",
        "Results ○●○",
        "04_human_profile_heatmap.png",
        [
            "Each image becomes a seven-emotion profile instead of one forced answer.",
            "Ambiguous images show more spread across emotion categories.",
            "This is why the analysis compares distributions rather than only dominant labels.",
        ],
    )

    fill_two_column_slide(
        prs.slides[15],
        "Dominant-emotion matrices show where humans and models collide",
        "Results ○●○",
        "10_dominant_emotion_matrices.png",
        [
            f"FER never selected disgust as dominant; DeepFace selected it only {n['df_disgust_dom']} times.",
            "Some human categories are repeatedly mapped to different AI categories.",
            "These collisions are more defensible to discuss than saying a model was simply wrong.",
        ],
    )

    fill_two_column_slide(
        prs.slides[16],
        "Agreement varies by intended emotion category",
        "Results ○●○",
        "11_agreement_by_assigned_category.png",
        [
            "Some emotion categories create more agreement than others.",
            "The assigned category is used only for grouping, not as certified ground truth.",
            "Low-agreement categories are useful evidence of ambiguity.",
        ],
    )

    fill_text_slide(
        prs.slides[17],
        "Modified images test stability under visual change",
        "Modified Images overview ●",
        [
            "Each original image was compared with blurred, greyscale, and low-resolution versions.",
            "The goal is to see whether emotion profiles stay stable when image information changes.",
            "This section separates model robustness, human interpretation, and ambiguity.",
        ],
    )

    fill_two_column_slide(
        prs.slides[18],
        "Blur affects DeepFace more than greyscale",
        "Modified Images ●○",
        "08_modification_robustness_full.png",
        [
            "Greyscale barely changes FER and DeepFace, likely because their preprocessing reduces color dependence.",
            f"DeepFace robustness is lower for blur ({n['df_blur']:.2f}) than greyscale ({n['df_grey']:.2f}).",
            "Human ratings move more across image variants than the AI models.",
        ],
    )

    fill_text_slide(
        prs.slides[19],
        "Greyscale makes human responses more mixed",
        "Modified Images ○●",
        [
            "The greyscale result should be described carefully.",
            f"In greyscale images, human dominant profiles included {n['human_grey_sad']} sad, {n['human_grey_neutral']} neutral, and {n['human_grey_disgust']} disgust cases.",
            "So the stronger claim is not simply “humans became worse.”",
            "Instead, greyscale appears to increase uncertainty around darker or negative-looking expressions.",
        ],
    )

    fill_two_column_slide(
        prs.slides[20],
        "Ambiguous images lower agreement and increase uncertainty",
        "Modified Images ○●",
        "07_ambiguity_panels.png",
        [
            "Ambiguous images produce higher human entropy.",
            "Dominant-emotion agreement drops on ambiguous images, especially for DeepFace.",
            "This supports the study’s focus on disagreement rather than correctness.",
        ],
    )

    fill_text_slide(
        prs.slides[21],
        "MediaPipe links emotion scores to facial movement features",
        "MediaPipe overview ●",
        [
            "MediaPipe extracts facial blendshape features, such as eye widening and mouth smile.",
            "These features are not emotion labels by themselves.",
            "We compare them with human, FER, and DeepFace emotion-score profiles.",
        ],
    )

    fill_two_column_slide(
        prs.slides[22],
        "Human ratings connect to intuitive facial features",
        "MediaPipe ●○○",
        "13a_mediapipe_top_associations_human_pooled.png",
        [
            "Eye widening is strongly associated with human surprise scores.",
            "Mouth smile features are associated with human happiness scores.",
            "This makes the human pooled ratings interpretable through visible facial movement.",
        ],
    )

    fill_two_column_slide(
        prs.slides[23],
        "FER associations emphasize classic expression cues",
        "MediaPipe ○●○",
        "13b_mediapipe_top_associations_fer.png",
        [
            "FER surprise is strongly associated with eye widening and jaw opening.",
            "FER happiness is associated with left and right mouth-smile features.",
            "This makes FER relatively interpretable through MediaPipe blendshapes.",
        ],
    )

    fill_two_column_slide(
        prs.slides[24],
        "DeepFace uses facial cues but aligns less with humans overall",
        "MediaPipe ○○●",
        "13c_mediapipe_top_associations_deepface.png",
        [
            "DeepFace still shows associations with visible facial features.",
            "However, its overall human similarity is lower than FER in this dataset.",
            "This suggests different models can read similar cues but weight them differently.",
        ],
    )

    # Keep the original thank-you slide style but fix missing spacing.
    for shape in prs.slides[25].shapes:
        text = getattr(shape, "text", "")
        if "THANK YOU" in text.upper():
            shape.text = "THANK YOU FOR LISTENING!\nANY QUESTIONS?"

    prs.save(str(OUT))
    print(f"Saved {OUT.resolve()}")


if __name__ == "__main__":
    main()
