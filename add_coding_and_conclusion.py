"""Add coding-depth and conclusion slides to the latest edited deck.

The script preserves the user's current deck and appends three normal content
slides before the final thank-you slide.
"""

from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


INPUT = Path("presentation/AI_LAB_expai_presentation_ACSAI_2025_2026_project_paragraphs.pptx")
OUTPUT = Path("presentation/AI_LAB_expai_presentation_ACSAI_2025_2026_with_coding_conclusion.pptx")

BLACK = RGBColor(0, 0, 0)


NEW_SLIDES = [
    (
        "The scripts made the data collection reproducible",
        "Coding ●○",
        (
            "The project used separate Python scripts for each source of data. The human-rating script presented images "
            "to participants in batches and saved their responses into CSV files. FER, DeepFace, and MediaPipe were then "
            "run with their own scripts so each system could process the same image versions in a repeatable way."
        ),
    ),
    (
        "The analysis code converted different outputs into one comparison format",
        "Coding ○●",
        (
            "The main analysis script cleaned the CSV files, matched rows by image ID and image version, and converted "
            "human and AI outputs into seven-emotion vectors. This made it possible to compare humans, FER, and DeepFace "
            "using the same measurements instead of treating each file separately."
        ),
    ),
    (
        "Conclusion: the project studies interpretation, not absolute correctness",
        "Conclusion ●",
        (
            "The main conclusion is that humans and AI systems do not always interpret facial expressions in the same way. "
            "FER aligned more closely with the participant ratings than DeepFace in this dataset, while modified and ambiguous "
            "images revealed where interpretation becomes less stable. MediaPipe helped connect these emotion scores to visible "
            "facial features, giving the project a stronger explanation layer."
        ),
    ),
]


def clone_shape_to_slide(shape, slide):
    new_element = copy.deepcopy(shape.element)
    slide.shapes._spTree.insert_element_before(new_element, "p:extLst")


def set_text_preserve_style(shape, text):
    nodes = shape.element.xpath(".//a:t")
    if not nodes:
        shape.text = text
        return
    nodes[0].text = text
    for node in nodes[1:]:
        node.text = ""


def delete_shape(shape):
    shape.element.getparent().remove(shape.element)


def clear_slide(slide):
    for shape in list(slide.shapes):
        delete_shape(shape)


def add_paragraph(slide, text):
    box = slide.shapes.add_textbox(Inches(1.9), Inches(1.75), Inches(12.2), Inches(4.6))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(24)
    p.font.color.rgb = BLACK
    p.line_spacing = 1.0
    p.space_before = Pt(11.91)
    p.space_after = Pt(9.92)


def build_content_slide(slide, chrome_slide, title, section, body):
    clear_slide(slide)
    # Copy the same template chrome used by the user's method pages:
    # footer bar, slide number/footer, logo, title and section label.
    for shape in list(chrome_slide.shapes)[:6]:
        clone_shape_to_slide(shape, slide)
    set_text_preserve_style(slide.shapes[4], title)
    set_text_preserve_style(slide.shapes[5], section)
    add_paragraph(slide, body)


def move_last_slide_before(prs, before_index_zero_based):
    """Move the last slide in the deck before the given zero-based index."""
    sld_id_lst = prs.slides._sldIdLst
    last = sld_id_lst[-1]
    sld_id_lst.remove(last)
    sld_id_lst.insert(before_index_zero_based, last)


def main():
    prs = Presentation(str(INPUT))
    chrome_slide = prs.slides[9]  # current methodology/coding slide

    # Insert new slides before the final thank-you slide, preserving everything
    # already in the user's deck.
    insert_at = len(prs.slides) - 1
    for title, section, body in NEW_SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        build_content_slide(slide, chrome_slide, title, section, body)
        move_last_slide_before(prs, insert_at)
        insert_at += 1

    prs.save(str(OUTPUT))
    print(f"Saved {OUTPUT.resolve()}")


if __name__ == "__main__":
    main()
