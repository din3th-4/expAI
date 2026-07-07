"""Rewrite generated presentation text as project-focused paragraphs.

This keeps the template chrome and graph images, but replaces the generated
bullet-style text with shorter paragraph prose using the same paragraph size and
spacing pattern as the user's Methodology pages.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt


INPUT = Path("presentation/AI_LAB_expai_presentation_ACSAI_2025_2026_no_red_backgrounds.pptx")
OUTPUT = Path("presentation/AI_LAB_expai_presentation_ACSAI_2025_2026_project_paragraphs.pptx")

BLACK = RGBColor(0, 0, 0)

PARAGRAPHS = {
    "Results roadmap": (
        "This part of the presentation moves from the collected data to the main findings of the project. "
        "Since the images were web-selected and not taken from a certified emotion dataset, the results are "
        "presented as a comparison of perception rather than a test of who is correct."
    ),
    "The dataset is usable, but not equally complete": (
        "Before comparing humans and AI models, we first checked whether each data source covered the same image set. "
        "The study combines participant ratings, FER, DeepFace, and MediaPipe outputs, so this slide shows how complete "
        "each source is and where missing detections may affect the comparison."
    ),
    "Agreement replaces accuracy in this study": (
        "Because the labels are researcher-assigned rather than certified ground truth, accuracy is not the right word "
        "for the main comparison. Instead, the project measures how closely each AI model behaves compared with the "
        "human pooled emotion profiles."
    ),
    "Average emotion profiles reveal model bias patterns": (
        "This graph compares the overall emotion tendencies of humans, FER, and DeepFace. It shows that the systems do "
        "not distribute emotion scores in the same way, which is important because the project is about patterns of "
        "agreement and disagreement, not just one final label."
    ),
    "Human ratings are multi-dimensional, not single labels": (
        "The human data was treated as a seven-emotion profile for each image. This is important because participants "
        "often saw more than one possible emotion in the same face, especially for ambiguous images."
    ),
    "Dominant-emotion matrices show where humans and models collide": (
        "This slide shows where the strongest human response and the strongest AI response matched or collided. These "
        "collisions are useful because they reveal the kinds of expressions that humans and AI systems interpret differently."
    ),
    "Agreement varies by intended emotion category": (
        "The researcher-assigned emotion category is used here only as a way to group the images. Some groups produced "
        "more agreement between humans and AI, while others created more disagreement and ambiguity."
    ),
    "Modified images test stability under visual change": (
        "After comparing the original images, the project tests whether emotion perception stays stable when the same "
        "faces are blurred, converted to greyscale, or reduced in resolution. This helps show whether humans and models "
        "depend on the same visual information."
    ),
    "Blur affects DeepFace more than greyscale": (
        "This graph compares each modified image with its original version. In this dataset, greyscale changes the AI "
        "outputs very little, while blur creates a larger shift for DeepFace."
    ),
    "Greyscale makes human responses more mixed": (
        "The greyscale images are interesting because they do not simply make people choose one emotion more often. "
        "Instead, they make responses more mixed around emotions such as sadness, neutrality, and disgust."
    ),
    "Ambiguous images lower agreement and increase uncertainty": (
        "The ambiguous images support the main idea of the project: emotion recognition is not always clear. When the "
        "image itself is harder to interpret, both human agreement and human-AI agreement become less stable."
    ),
    "MediaPipe links emotion scores to facial movement features": (
        "MediaPipe adds another layer to the study by extracting facial movement features from the same images. These "
        "features are not emotion labels, but they help explain which visible facial cues are related to human and AI responses."
    ),
    "Human ratings connect to intuitive facial features": (
        "For human pooled responses, the strongest MediaPipe associations are easy to interpret. Eye widening connects "
        "with surprise, while mouth-smile features connect with happiness."
    ),
    "FER associations emphasize classic expression cues": (
        "FER shows a similar connection to classic facial-expression cues. Surprise relates strongly to eye widening and "
        "jaw opening, while happiness relates to mouth-smile features."
    ),
    "DeepFace uses facial cues but aligns less with humans overall": (
        "DeepFace also responds to visible facial features, but its emotion profiles align less with the human pooled "
        "ratings than FER in this dataset. This suggests that the models may use similar facial information but weight it differently."
    ),
}


def delete_shape(shape):
    shape.element.getparent().remove(shape.element)


def slide_title(slide):
    for shape in slide.shapes:
        text = getattr(shape, "text", "").strip()
        if text in PARAGRAPHS:
            return text
    return None


def body_area(slide):
    """Return paragraph placement based on whether a graph image is present."""
    # Template/logo images are near the top-left; content graph images are
    # larger and start below the title area.
    content_pictures = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.top > Inches(1.0)
    ]
    if content_pictures:
        return Inches(8.45), Inches(1.85), Inches(6.35), Inches(4.6), Pt(22)
    return Inches(1.9), Inches(1.7), Inches(12.2), Inches(4.5), Pt(24)


def remove_generated_textboxes(slide):
    # Preserve template chrome/title/section and all pictures. Remove generated
    # text-containing shapes in the content area.
    for shape in list(slide.shapes):
        has_text = bool(getattr(shape, "text", "").strip())
        if (
            has_text
            and shape.shape_type != MSO_SHAPE_TYPE.PICTURE
            and shape.top > Inches(1.0)
        ):
            delete_shape(shape)


def add_project_paragraph(slide, text):
    left, top, width, height, size = body_area(slide)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = size
    p.font.color.rgb = BLACK
    p.line_spacing = 1.0
    p.space_before = Pt(11.91)
    p.space_after = Pt(9.92)


def main():
    prs = Presentation(str(INPUT))
    edited = 0

    for slide in prs.slides:
        title = slide_title(slide)
        if not title:
            continue
        remove_generated_textboxes(slide)
        add_project_paragraph(slide, PARAGRAPHS[title])
        edited += 1

    prs.save(str(OUTPUT))
    print(f"Rewrote {edited} generated slides")
    print(f"Saved {OUTPUT.resolve()}")


if __name__ == "__main__":
    main()
